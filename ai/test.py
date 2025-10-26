import os
import fitz  # PyMuPDF
import pptx
import json
import re
from typing import List, Dict, Any
from datetime import datetime
import hashlib

class KnowledgeBaseCreator:
    """
    Enhanced knowledge base creator for RAG systems with optimized chunking,
    metadata extraction, and document processing.
    """
    
    def __init__(self, output_dir_images='extracted_images', 
                 chunk_size=800, chunk_overlap=100):
        """
        Initialize the knowledge base creator.
        
        Args:
            output_dir_images: Directory to save extracted images
            chunk_size: Target size for text chunks (characters)
            chunk_overlap: Overlap between chunks to maintain context
        """
        self.output_dir_images = os.path.join(os.getcwd(), output_dir_images)
        os.makedirs(self.output_dir_images, exist_ok=True)
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        
    def semantic_chunking(self, text: str, metadata: Dict) -> List[Dict]:
        """
        Performs semantic chunking with overlap to preserve context.
        Splits on natural boundaries like paragraphs and sentences.
        """
        # First, clean the text
        cleaned_text = re.sub(r'\s+', ' ', text).strip()
        
        # Split on paragraph boundaries first
        paragraphs = re.split(r'\n\s*\n', cleaned_text)
        
        chunks = []
        current_chunk = ""
        current_length = 0
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
                
            para_length = len(para)
            
            # If paragraph itself is too long, split by sentences
            if para_length > self.chunk_size:
                sentences = re.split(r'(?<=[.!?])\s+', para)
                for sentence in sentences:
                    if current_length + len(sentence) > self.chunk_size and current_chunk:
                        # Save current chunk
                        if len(current_chunk.strip()) > 100:  # Minimum chunk size
                            chunks.append(self._create_chunk(current_chunk, metadata))
                        
                        # Start new chunk with overlap
                        overlap_text = self._get_overlap(current_chunk)
                        current_chunk = overlap_text + " " + sentence
                        current_length = len(current_chunk)
                    else:
                        current_chunk += " " + sentence
                        current_length += len(sentence) + 1
            else:
                # Add paragraph to current chunk
                if current_length + para_length > self.chunk_size and current_chunk:
                    if len(current_chunk.strip()) > 100:
                        chunks.append(self._create_chunk(current_chunk, metadata))
                    
                    overlap_text = self._get_overlap(current_chunk)
                    current_chunk = overlap_text + " " + para
                    current_length = len(current_chunk)
                else:
                    if current_chunk:
                        current_chunk += "\n\n" + para
                    else:
                        current_chunk = para
                    current_length += para_length + 2
        
        # Add remaining chunk
        if current_chunk.strip() and len(current_chunk.strip()) > 100:
            chunks.append(self._create_chunk(current_chunk, metadata))
        
        return chunks
    
    def _get_overlap(self, text: str) -> str:
        """Extract overlap text from end of chunk."""
        if len(text) <= self.chunk_overlap:
            return text
        return text[-self.chunk_overlap:]
    
    def _create_chunk(self, content: str, metadata: Dict) -> Dict:
        """Create a chunk with enhanced metadata."""
        chunk_hash = hashlib.md5(content.encode()).hexdigest()[:8]
        return {
            "content": content.strip(),
            "metadata": metadata.copy(),
            "chunk_hash": chunk_hash,
            "char_count": len(content.strip()),
            "word_count": len(content.strip().split())
        }
    
    def extract_metadata(self, file_path: str, filename: str, 
                        page_or_slide_num: int, doc_type: str) -> Dict:
        """
        Extract comprehensive metadata for better retrieval.
        """
        file_stats = os.stat(file_path)
        
        metadata = {
            "source_file": filename,
            "source_type": doc_type,
            "page_number": page_or_slide_num,
            "source_reference": f"{filename}-{doc_type}-{page_or_slide_num}",
            "file_size_bytes": file_stats.st_size,
            "created_date": datetime.fromtimestamp(file_stats.st_ctime).isoformat(),
            "modified_date": datetime.fromtimestamp(file_stats.st_mtime).isoformat(),
            "extraction_date": datetime.now().isoformat()
        }
        
        return metadata
    
    def process_pdf(self, file_path: str, filename_base: str) -> List[Dict]:
        """
        Enhanced PDF processing with better text extraction and metadata.
        """
        document_chunks = []
        
        try:
            doc = fitz.open(file_path)
            
            for page_num, page in enumerate(doc):
                # Extract metadata
                metadata = self.extract_metadata(
                    file_path, filename_base, page_num + 1, "pdf-page"
                )
                
                # Get text blocks in reading order
                blocks = page.get_text("blocks", sort=True)
                page_text = ""
                
                for block in blocks:
                    block_type = block[6]
                    
                    if block_type == 0:  # Text block
                        text = block[4].strip()
                        if text:
                            page_text += text + "\n\n"
                    
                    elif block_type == 1:  # Image block
                        # Skip full-page images
                        image_bbox = fitz.Rect(block[:4])
                        page_rect = page.rect
                        
                        if (image_bbox.width > page_rect.width * 0.95 and 
                            image_bbox.height > page_rect.height * 0.95):
                            continue
                        
                        # Extract and save image
                        try:
                            xref = block[5]
                            image_data = doc.extract_image(xref)
                            image_bytes = image_data["image"]
                            image_ext = image_data["ext"]
                            
                            image_filename = f"{filename_base}_p{page_num+1}_img_{xref}.{image_ext}"
                            image_path = os.path.join(self.output_dir_images, image_filename)
                            
                            with open(image_path, "wb") as img_file:
                                img_file.write(image_bytes)
                            
                            # Add image reference with metadata
                            image_metadata = metadata.copy()
                            image_metadata["content_type"] = "image"
                            image_metadata["image_filename"] = image_filename
                            
                            document_chunks.append({
                                "content": f"[IMAGE: {image_filename}]",
                                "metadata": image_metadata,
                                "content_type": "image"
                            })
                        except Exception as e:
                            print(f"Error extracting image on page {page_num+1}: {e}")
                
                # Chunk the page text semantically
                if page_text.strip():
                    text_metadata = metadata.copy()
                    text_metadata["content_type"] = "text"
                    chunks = self.semantic_chunking(page_text, text_metadata)
                    document_chunks.extend(chunks)
            
            doc.close()
            
        except Exception as e:
            print(f"Error processing PDF {file_path}: {e}")
        
        return document_chunks
    
    def process_ppt(self, file_path: str, filename_base: str) -> List[Dict]:
        """
        Enhanced PPT processing with better structure preservation.
        """
        document_chunks = []
        
        try:
            presentation = pptx.Presentation(file_path)
            
            for slide_num, slide in enumerate(presentation.slides):
                # Extract metadata
                metadata = self.extract_metadata(
                    file_path, filename_base, slide_num + 1, "ppt-slide"
                )
                
                slide_text = ""
                
                # Extract title if present
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        metadata["slide_title"] = title_text
                        slide_text += f"# {title_text}\n\n"
                
                # Process all shapes
                for shape in slide.shapes:
                    # Extract text
                    if hasattr(shape, "text") and shape.has_text_frame:
                        text = shape.text.strip()
                        if text and text != metadata.get("slide_title", ""):
                            slide_text += text + "\n\n"
                    
                    # Extract images
                    if hasattr(shape, 'image'):
                        try:
                            image_bytes = shape.image.blob
                            image_ext = shape.image.ext
                            image_filename = f"{filename_base}_s{slide_num+1}_img_{shape.shape_id}.{image_ext}"
                            image_path = os.path.join(self.output_dir_images, image_filename)
                            
                            with open(image_path, "wb") as img_file:
                                img_file.write(image_bytes)
                            
                            # Add image reference
                            image_metadata = metadata.copy()
                            image_metadata["content_type"] = "image"
                            image_metadata["image_filename"] = image_filename
                            
                            document_chunks.append({
                                "content": f"[IMAGE: {image_filename}]",
                                "metadata": image_metadata,
                                "content_type": "image"
                            })
                        except Exception as e:
                            print(f"Error extracting image from slide {slide_num+1}: {e}")
                
                # Chunk the slide text
                if slide_text.strip():
                    text_metadata = metadata.copy()
                    text_metadata["content_type"] = "text"
                    chunks = self.semantic_chunking(slide_text, text_metadata)
                    document_chunks.extend(chunks)
        
        except Exception as e:
            print(f"Error processing PPT {file_path}: {e}")
        
        return document_chunks
    
    def create_knowledge_base(self, folder_path: str) -> List[Dict]:
        """
        Main function to process all files and create knowledge base.
        """
        all_chunks = []
        
        supported_files = [f for f in os.listdir(folder_path) 
                          if f.endswith(('.pdf', '.pptx'))]
        
        print(f"Found {len(supported_files)} supported files to process.")
        
        for filename in supported_files:
            file_path = os.path.join(folder_path, filename)
            filename_base = os.path.splitext(filename)[0]
            
            print(f"\nProcessing: {filename}")
            
            if filename.endswith('.pdf'):
                chunks = self.process_pdf(file_path, filename_base)
                all_chunks.extend(chunks)
                print(f"  ✓ Extracted {len(chunks)} chunks from PDF")
            
            elif filename.endswith('.pptx'):
                chunks = self.process_ppt(file_path, filename_base)
                all_chunks.extend(chunks)
                print(f"  ✓ Extracted {len(chunks)} chunks from PPT")
        
        # Add unique IDs to all chunks
        for i, chunk_data in enumerate(all_chunks):
            chunk_data["id"] = f"chunk_{i:05d}"
        
        return all_chunks
    
    def save_knowledge_base(self, knowledge_base: List[Dict], 
                           output_file: str = 'knowledge_base.json'):
        """
        Save knowledge base with statistics.
        """
        # Generate statistics
        stats = {
            "total_chunks": len(knowledge_base),
            "text_chunks": sum(1 for c in knowledge_base 
                             if c.get("content_type") != "image"),
            "image_references": sum(1 for c in knowledge_base 
                                  if c.get("content_type") == "image"),
            "total_characters": sum(c.get("char_count", 0) 
                                  for c in knowledge_base),
            "average_chunk_size": sum(c.get("char_count", 0) 
                                    for c in knowledge_base) / len(knowledge_base) 
                                    if knowledge_base else 0,
            "creation_date": datetime.now().isoformat()
        }
        
        output_data = {
            "metadata": stats,
            "chunks": knowledge_base
        }
        
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(output_data, f, indent=2, ensure_ascii=False)
        
        print(f"\n{'='*60}")
        print("Knowledge Base Creation Summary:")
        print(f"{'='*60}")
        print(f"Total chunks: {stats['total_chunks']}")
        print(f"Text chunks: {stats['text_chunks']}")
        print(f"Image references: {stats['image_references']}")
        print(f"Average chunk size: {stats['average_chunk_size']:.0f} characters")
        print(f"Total characters: {stats['total_characters']:,}")
        print(f"\nKnowledge base saved to: {output_file}")
        print(f"Images saved to: {self.output_dir_images}")
        print(f"{'='*60}")


if __name__ == "__main__":
    # Configuration
    DATA_FOLDER = 'D:\\NHITM\\BE\\BE Major Project\\code\\notes'  # Update this path
    OUTPUT_FILE = 'knowledge_base.json'
    
    # Optimal chunk size for RAG (600-1000 characters recommended)
    CHUNK_SIZE = 800
    CHUNK_OVERLAP = 100  # Maintains context between chunks
    
    if not os.path.isdir(DATA_FOLDER):
        print(f"Error: The folder '{DATA_FOLDER}' does not exist.")
    else:
        print("Starting Knowledge Base Creation for RAG System")
        print(f"Chunk Size: {CHUNK_SIZE} characters")
        print(f"Chunk Overlap: {CHUNK_OVERLAP} characters\n")
        
        # Create knowledge base
        creator = KnowledgeBaseCreator(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP
        )
        
        knowledge_base = creator.create_knowledge_base(DATA_FOLDER)
        creator.save_knowledge_base(knowledge_base, OUTPUT_FILE)
